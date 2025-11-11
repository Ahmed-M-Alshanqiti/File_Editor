# create_dummy_files.py
from pathlib import Path
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from fpdf import FPDF

base_dir = Path(__file__).resolve().parent / "dummy_files"
base_dir.mkdir(exist_ok=True)

# 1. DOCX
doc = Document()
doc.add_heading('Sample Document', 0)
doc.add_paragraph('This is a dummy Word document for testing LibreOffice conversion.')
doc.save(base_dir / 'sample_doc.docx')

# 2. XLSX
wb = Workbook()
ws = wb.active
ws['A1'] = "Sample Excel Data"
ws['A2'] = 12345
wb.save(base_dir / 'sample_sheet.xlsx')

# 3. PPTX
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Sample Presentation"
slide.placeholders[1].text = "This is a dummy slide for LibreOffice testing."
prs.save(base_dir / 'sample_presentation.pptx')

# 4. TXT
(base_dir / 'example_text.txt').write_text("Hello, this is a plain text file for upload testing.")

# 5. PDF
pdf = FPDF()
pdf.add_page()
pdf.set_font("Arial", size=12)
pdf.cell(200, 10, txt="This is a sample PDF file.", ln=True, align="C")
pdf.output(str(base_dir / "demo_pdf.pdf"))

# 6. JPG (simple placeholder)
(base_dir / 'image_example.jpg').write_bytes(b'\xff\xd8\xff\xe0' + b'\x00'*100 + b'\xff\xd9')

print(f"Dummy files created in: {base_dir}")
